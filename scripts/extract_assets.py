import base64
import io
import json
import mimetypes
import os
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook

TEXT_EXTENSIONS = {
    ".txt", ".md", ".json", ".js", ".ts", ".tsx", ".jsx", ".html", ".htm", ".css",
    ".py", ".java", ".c", ".cpp", ".cs", ".go", ".rs", ".php", ".rb", ".yml", ".yaml",
    ".xml", ".csv", ".tsv", ".ini", ".env", ".toml", ".sql", ".sh", ".ps1", ".log",
    ".bat", ".cmd", ".svg", ".mjs", ".vue"
}

CHUNK_SIZE = 1800
CHUNK_OVERLAP = 240
MAX_CONTENT_PREVIEW = 24000
MAX_FRAMES = 3


def chunk_text(text: str):
    text = text or ""
    if not text:
      return []

    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + CHUNK_SIZE)
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
    return chunks


def read_text_file(path: Path):
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except Exception:
            continue
    raise ValueError("Could not decode this text file.")


def extract_pdf(path: Path):
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append(f"Page {index + 1}\n{text}".strip())
    return "\n\n".join(pages).strip()


def extract_docx(path: Path):
    doc = Document(str(path))
    parts = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))

    return "\n".join(parts).strip()


def extract_pptx(path: Path):
    pres = Presentation(str(path))
    slides = []
    for index, slide in enumerate(pres.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                value = shape.text.strip()
                if value:
                    texts.append(value)
        if texts:
            slides.append(f"Slide {index + 1}\n" + "\n".join(texts))
    return "\n\n".join(slides).strip()


def extract_xlsx(path: Path):
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts).strip()


def image_to_base64(path: Path):
    with Image.open(path) as img:
        img = img.convert("RGB")
        img.thumbnail((1280, 1280))
        width, height = img.size
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=88)
        return base64.b64encode(buf.getvalue()).decode("utf-8"), width, height


def ffprobe_metadata(path: Path):
    result = subprocess.run(
        ["ffprobe", "-v", "error", "-print_format", "json", "-show_streams", "-show_format", str(path)],
        capture_output=True,
        text=True,
        check=True
    )
    return json.loads(result.stdout or "{}")


def extract_video_frames(path: Path):
    meta = ffprobe_metadata(path)
    duration = float(meta.get("format", {}).get("duration") or 0)
    video_stream = next((s for s in meta.get("streams", []) if s.get("codec_type") == "video"), {})
    width = int(video_stream.get("width") or 0)
    height = int(video_stream.get("height") or 0)

    timestamps = []
    if duration > 0:
        for ratio in (0.15, 0.5, 0.85):
            timestamps.append(max(0.0, duration * ratio))
    else:
        timestamps = [0.0]

    frames = []
    with tempfile.TemporaryDirectory() as tmpdir:
        for index, timestamp in enumerate(timestamps[:MAX_FRAMES]):
            out_path = Path(tmpdir) / f"frame-{index}.jpg"
            subprocess.run(
                [
                    "ffmpeg", "-y", "-ss", str(timestamp), "-i", str(path),
                    "-frames:v", "1", "-vf", "scale='min(960,iw)':-2", str(out_path)
                ],
                capture_output=True,
                check=True
            )
            if out_path.exists():
                frames.append(base64.b64encode(out_path.read_bytes()).decode("utf-8"))

    return {
        "duration": duration,
        "width": width,
        "height": height,
        "frames": frames
    }


def make_payload(path_str: str):
    path = Path(path_str)
    if not path.exists():
        raise ValueError("File not found.")
    if not path.is_file():
        raise ValueError("Only files are supported.")

    mime_type, _ = mimetypes.guess_type(str(path))
    ext = path.suffix.lower()
    stat = path.stat()

    payload = {
        "name": path.name,
        "path": str(path),
        "size": stat.st_size,
        "mimeType": mime_type or "application/octet-stream",
        "content": "",
        "chunks": [],
        "truncated": False,
        "charCount": 0,
        "mediaKind": "text"
    }

    if ext in TEXT_EXTENSIONS:
        text = read_text_file(path)
        payload["content"] = text[:MAX_CONTENT_PREVIEW]
        payload["chunks"] = chunk_text(text)
        payload["truncated"] = len(text) > MAX_CONTENT_PREVIEW
        payload["charCount"] = len(text)
        return payload

    if ext == ".pdf":
        text = extract_pdf(path)
        payload["content"] = text[:MAX_CONTENT_PREVIEW]
        payload["chunks"] = chunk_text(text)
        payload["truncated"] = len(text) > MAX_CONTENT_PREVIEW
        payload["charCount"] = len(text)
        return payload

    if ext == ".docx":
        text = extract_docx(path)
        payload["content"] = text[:MAX_CONTENT_PREVIEW]
        payload["chunks"] = chunk_text(text)
        payload["truncated"] = len(text) > MAX_CONTENT_PREVIEW
        payload["charCount"] = len(text)
        return payload

    if ext == ".pptx":
        text = extract_pptx(path)
        payload["content"] = text[:MAX_CONTENT_PREVIEW]
        payload["chunks"] = chunk_text(text)
        payload["truncated"] = len(text) > MAX_CONTENT_PREVIEW
        payload["charCount"] = len(text)
        return payload

    if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        text = extract_xlsx(path)
        payload["content"] = text[:MAX_CONTENT_PREVIEW]
        payload["chunks"] = chunk_text(text)
        payload["truncated"] = len(text) > MAX_CONTENT_PREVIEW
        payload["charCount"] = len(text)
        return payload

    if (mime_type or "").startswith("image/"):
        image_b64, width, height = image_to_base64(path)
        summary = f"Image file {path.name} ({width}x{height})."
        payload["mediaKind"] = "image"
        payload["imageBase64"] = image_b64
        payload["width"] = width
        payload["height"] = height
        payload["content"] = summary
        payload["chunks"] = [summary]
        payload["charCount"] = len(summary)
        return payload

    if (mime_type or "").startswith("video/") or ext in {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v"}:
        video = extract_video_frames(path)
        duration = round(video["duration"], 2)
        summary = f"Video file {path.name} ({video['width']}x{video['height']}, {duration}s)."
        payload["mediaKind"] = "video"
        payload["videoFramesBase64"] = video["frames"]
        payload["width"] = video["width"]
        payload["height"] = video["height"]
        payload["duration"] = duration
        payload["content"] = summary
        payload["chunks"] = [summary]
        payload["charCount"] = len(summary)
        return payload

    raise ValueError("This file format is not supported yet.")


def main():
    try:
        if len(sys.argv) > 1:
            paths = sys.argv[1:]
        else:
            raw = sys.stdin.read().lstrip("\ufeff").strip()
            data = json.loads(raw or "{}")
            paths = data.get("paths") or []

        out = {"files": [], "errors": []}

        for path_str in paths:
            try:
                out["files"].append(make_payload(path_str))
            except Exception as exc:
                out["errors"].append({
                    "path": str(path_str),
                    "name": os.path.basename(str(path_str)),
                    "error": str(exc)
                })

        sys.stdout.write(json.dumps(out))
    except Exception as exc:
        sys.stdout.write(json.dumps({"files": [], "errors": [{"path": "", "name": "", "error": str(exc)}]}))


if __name__ == "__main__":
    main()
